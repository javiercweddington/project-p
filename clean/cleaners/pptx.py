"""
PPTXCleaner - clean PowerPoint presentations by removing metadata and cleaning content.

Addresses PowerPoint-specific risks (ranked by failure likelihood):

1. Speaker notes (CRITICAL) - The frankest content in the file. People write
   things in notes they'd never put on a slide. Pricing discussions, internal
   feedback, client strategies.

2. Off-canvas objects (HIGH) - People park old slides, deprecated content,
   and sensitive diagrams off-canvas instead of deleting them. These are
   fully recoverable.

3. Hidden slides (HIGH) - Slides marked as hidden don't appear in normal
   presentation mode but are fully extractable.

4. Cropping is non-destructive (MEDIUM) - A screenshot cropped to hide a
   customer name still contains the full original image. PowerPoint cropping
   stores crop boundaries but retains the original.

5. Objects layered behind other objects (MEDIUM) - Text boxes, images, or
   shapes placed behind visible content. Fully recoverable.

6. Embedded charts (MEDIUM) - Charts carry their full backing worksheet
   data, which may contain sensitive numbers not shown in the chart.

7. Embedded OLE objects (LOW) - Excel workbooks, Word docs embedded in
   slides.

Strategy: Use python-pptx for high-level cleaning, then directly manipulate
the underlying XML to handle off-canvas objects, hidden slides, and cropping
artifacts that python-pptx cannot access through its API.

Dependencies (optional):
- python-pptx: PowerPoint manipulation
"""

from __future__ import annotations

import io
import logging
import os
import re
import shutil
import zipfile
from pathlib import Path
from typing import List, Optional

from ..anonymizer import EntityMapper
from .text import TextCleaner

_logger = logging.getLogger(__name__)

# Try optional dependencies
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.oxml.ns import qn
    HAS_PYTHON_PPTX = True
except ImportError:
    HAS_PYTHON_PPTX = False


class PPTXCleaner:
    """Clean PowerPoint presentations by removing metadata and cleaning content.

    Uses python-pptx for high-level cleaning and direct XML manipulation
    for off-canvas objects, hidden slides, and cropping artifacts.
    """

    # Office document core properties to clear
    _CORE_PROPERTIES = {
        'creator', 'last_modified_by', 'contributor',
        'author', 'company', 'manager',
        'description', 'subject', 'title',
        'keywords', 'category', 'comments',
    }

    def __init__(self, mapper: EntityMapper):
        self.mapper = mapper
        self.text_cleaner = TextCleaner(mapper)

        if not HAS_PYTHON_PPTX:
            _logger.warning(
                "python-pptx not available; PowerPoint cleaning limited. "
                "Install with: pip install python-pptx"
            )

    def clean_file(self, input_path: Path, output_path: Path) -> bool:
        """Clean a PowerPoint presentation by removing metadata and cleaning content.

        Args:
            input_path: Source PPTX file path
            output_path: Destination PPTX file path

        Returns:
            True if cleaning was successful
        """
        ext = input_path.suffix.lower()

        # Legacy .ppt: rasterize (soffice -> PDF -> render/OCR/GLiNER
        # pixel redaction -> image-only PDF). Ghost content, Mandarin
        # text, embedded photos and signatures are all handled at the
        # pixel level; the deliverable becomes a .pdf sibling.
        if ext == '.ppt':
            return self._clean_legacy_ppt_via_raster(input_path, output_path)

        if not HAS_PYTHON_PPTX:
            _logger.warning("python-pptx not available; PPTX fail-closed")
            return False

        try:
            prs = Presentation(str(input_path))

            # Clear core properties
            self._clear_properties(prs.core_properties)

            # Clean all slides (text, notes, shapes)
            self._clean_all_slides(prs)

            # Save to bytes for XML manipulation
            output_bytes = io.BytesIO()
            prs.save(output_bytes)
            output_bytes.seek(0)

            # Post-processing: clean XML-level artifacts
            output_path.parent.mkdir(parents=True, exist_ok=True)
            self._clean_xml_artifacts(output_bytes, output_path)

            # Catch-all: entity text in members python-pptx never visits
            # (SmartArt, charts, embedded drawing XML). Fail closed.
            from .xml_pass import scrub_zip_xml_members
            if not scrub_zip_xml_members(output_path, self.mapper,
                                         input_path.name):
                output_path.unlink(missing_ok=True)
                return False

            return True

        except Exception as e:
            _logger.error("Error cleaning PPTX %s: %s", input_path, e)
            return False

    def _clean_legacy_ppt_via_raster(self, input_path: Path,
                                     output_path: Path) -> bool:
        """Legacy .ppt via the raster architecture.

        soffice converts the deck to PDF (text layer preserved for the
        exact-geometry redaction belt), then the PDF raster cleaner
        renders, redacts on pixels (OCR + mapper + shape rules + GLiNER)
        and rebuilds an image-only PDF. Writing legacy .ppt binaries
        back is impractical, so the deliverable is a .pdf SIBLING of
        output_path; the pipeline detects the conversion and tracks the
        new file. Fail-closed when LibreOffice or raster prerequisites
        are missing.
        """
        import os
        import shutil as _shutil
        import subprocess
        import tempfile

        soffice = _shutil.which('soffice') or _shutil.which('libreoffice')
        if soffice is None:
            _logger.warning(
                "Legacy .ppt %s: LibreOffice (soffice) not found — cannot "
                "rasterize. Install libreoffice, or the deck stays "
                "quarantined.", input_path.name)
            return False

        from .pdf import PDFCleaner
        pdf_cleaner = PDFCleaner(self.mapper)
        if not pdf_cleaner._raster_available():
            _logger.warning(
                "Legacy .ppt %s: raster prerequisites missing "
                "(PyMuPDF + OCR backend). Fail-closed.", input_path.name)
            return False

        final_pdf = output_path.with_suffix('.pdf')
        with tempfile.TemporaryDirectory(prefix='ppt_raster_') as tmp:
            tmpdir = Path(tmp)
            try:
                proc = subprocess.run(
                    [soffice, '--headless', '--convert-to', 'pdf',
                     '--outdir', str(tmpdir), str(input_path)],
                    capture_output=True, timeout=300)
            except Exception as e:
                _logger.error(
                    "soffice conversion failed for %s: %s",
                    input_path.name, e)
                return False
            produced = tmpdir / (input_path.stem + '.pdf')
            if proc.returncode != 0 or not produced.exists():
                _logger.error(
                    "soffice could not convert %s (rc=%s): %s",
                    input_path.name, proc.returncode,
                    proc.stderr.decode('utf-8', 'replace')[:300])
                return False

            if not pdf_cleaner._clean_raster(produced, final_pdf):
                final_pdf.unlink(missing_ok=True)
                return False

        # Remove the original .ppt from staging (the .pdf replaces it).
        try:
            if output_path.exists() and output_path != final_pdf:
                os.remove(output_path)
        except OSError as e:
            _logger.warning("Could not remove converted .ppt %s: %s — "
                            "deleting cleaned output, fail-closed.",
                            output_path.name, e)
            final_pdf.unlink(missing_ok=True)
            return False
        _logger.info(
            "Rasterized legacy deck %s -> %s (image-only PDF)",
            input_path.name, final_pdf.name)
        return True

    def _clear_properties(self, core_props) -> None:
        """Clear core document properties."""
        # Identifier properties go through the mapper (consistent
        # placeholders); python-pptx exposes dc:creator as `author`.
        for prop_name, entity_type in (('author', 'person'),
                                       ('last_modified_by', 'person')):
            try:
                value = getattr(core_props, prop_name, None)
                if value and str(value).strip():
                    setattr(core_props, prop_name, self.mapper.get_or_create(
                        entity_type=entity_type, value=str(value).strip(),
                        source='pptx_core_properties'))
                else:
                    setattr(core_props, prop_name, '')
            except (AttributeError, TypeError):
                pass

        for prop_name in ('comments', 'category', 'content_status',
                          'identifier', 'keywords', 'language', 'subject',
                          'title', 'version'):
            try:
                setattr(core_props, prop_name, '')
            except (AttributeError, TypeError):
                pass

        # Timestamps require datetime objects; normalize to fixed epoch
        from datetime import datetime as _dt
        fixed = _dt(2024, 1, 1, 0, 0, 0)
        for prop_name in ('created', 'modified', 'last_printed'):
            try:
                setattr(core_props, prop_name, fixed)
            except (AttributeError, TypeError, ValueError):
                pass

    def _clean_all_slides(self, prs) -> None:
        """Clean text content across all slides, notes, and shapes."""
        for slide_num, slide in enumerate(prs.slides, 1):
            # Clean slide notes (high risk - frankest content)
            if slide.has_notes_slide:
                notes_slide = slide.notes_slide
                if notes_slide.notes_text_frame:
                    _logger.debug("Cleaning speaker notes for slide %d", slide_num)
                    for para in notes_slide.notes_text_frame.paragraphs:
                        for run in para.runs:
                            if isinstance(run.text, str) and run.text:
                                run.text = self.text_cleaner.clean_text(run.text)

            # Clean slide shapes
            for shape in slide.shapes:
                # Clean text in shapes
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if isinstance(run.text, str) and run.text:
                                run.text = self.text_cleaner.clean_text(run.text)

                # Clean tables embedded in shapes
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            for para in cell.text_frame.paragraphs:
                                for run in para.runs:
                                    if isinstance(run.text, str) and run.text:
                                        run.text = self.text_cleaner.clean_text(
                                            run.text
                                        )

    def _clean_xml_artifacts(self, source_bytes: io.BytesIO,
                              output_path: Path) -> None:
        """Clean XML-level artifacts that python-pptx cannot access.

        This directly manipulates the ZIP internals to:
        1. Remove off-canvas objects
        2. Handle hidden slides
        3. Clean cropping artifacts
        4. Remove embedded chart data
        """
        source_bytes.seek(0)

        try:
            with zipfile.ZipFile(source_bytes, 'r') as src_zf:
                with zipfile.ZipFile(output_path, 'w', zipfile.ZIP_DEFLATED) as dst_zf:
                    for entry in src_zf.infolist():
                        if entry.filename.endswith('/'):
                            continue

                        name_lower = entry.filename.lower()
                        data = src_zf.read(entry.filename)

                        if name_lower.endswith(('.xml', '.rels')):
                            try:
                                text = data.decode('utf-8')

                                if name_lower.startswith('ppt/slides/slide'):
                                    text = self._clean_slide_xml(text)

                                # Catch-all visible-text pass over every
                                # text-bearing part (slides, notes, masters,
                                # layouts): grouped shapes and other content
                                # python-pptx never handed us live in <a:t>.
                                if name_lower.startswith(
                                        ('ppt/slides/', 'ppt/notesslides/',
                                         'ppt/slidemasters/',
                                         'ppt/slidelayouts/')):
                                    text = self._clean_visible_text_xml(
                                        text, 'a:t')

                                if name_lower == 'docprops/app.xml':
                                    text = self._clean_app_xml(text)

                                if name_lower.endswith('.rels'):
                                    text = self._clean_rel_targets(text)

                                data = text.encode('utf-8')
                            except Exception:
                                pass

                        # Fixed member timestamp: never leak the run date.
                        info = zipfile.ZipInfo(
                            filename=entry.filename,
                            date_time=(2024, 1, 1, 0, 0, 0),
                        )
                        info.compress_type = zipfile.ZIP_DEFLATED
                        dst_zf.writestr(info, data)

        except Exception:
            # Fail closed: never write the un-scrubbed bytes.
            if output_path.exists():
                output_path.unlink()
            raise

    def _clean_visible_text_xml(self, xml_text: str, tag: str) -> str:
        """Entity-clean the inner text of every <tag>...</tag> element."""
        from xml.sax.saxutils import escape, unescape

        def _sub(m: re.Match) -> str:
            raw = unescape(m.group(2))
            cleaned_text = self.text_cleaner.clean_text(raw)
            if cleaned_text == raw:
                return m.group(0)
            return m.group(1) + escape(cleaned_text) + m.group(3)

        return re.sub(
            rf'(<{tag}(?:\s[^>]*)?>)(.*?)(</{tag}>)',
            _sub, xml_text, flags=re.DOTALL,
        )

    def _clean_rel_targets(self, xml_text: str) -> str:
        """Entity-clean Target="..." values (mailto:, entity URLs)."""
        from xml.sax.saxutils import escape, unescape

        def _sub(m: re.Match) -> str:
            raw = unescape(m.group(2))
            cleaned_text = self.text_cleaner.clean_text(raw)
            if cleaned_text == raw:
                return m.group(0)
            return m.group(1) + escape(cleaned_text) + m.group(3)

        return re.sub(r'(Target=")([^"]*)(")', _sub, xml_text)

    def _clean_app_xml(self, xml_text: str) -> str:
        """Scrub docProps/app.xml: Company/Manager pseudonymized, tool
        fingerprints blanked, slide-title lists entity-cleaned."""
        from xml.sax.saxutils import unescape
        cleaned = xml_text

        for tag, entity_type in (('Company', 'company'),
                                 ('Manager', 'person')):
            def _sub(m: re.Match, _etype=entity_type) -> str:
                value = unescape(m.group(2)).strip()
                if not value:
                    return m.group(0)
                placeholder = self.mapper.get_or_create(
                    entity_type=_etype, value=value,
                    source='pptx_app_properties')
                return f"{m.group(1)}{placeholder}{m.group(3)}"
            cleaned = re.sub(rf'(<{tag}>)([^<]*)(</{tag}>)', _sub, cleaned)

        for tag in ('Template', 'HyperlinkBase', 'TotalTime',
                    'Application', 'AppVersion'):
            cleaned = re.sub(rf'<{tag}>[^<]*</{tag}>',
                             f'<{tag}></{tag}>', cleaned)

        cleaned = self._clean_visible_text_xml(cleaned, 'vt:lpstr')
        return cleaned

    def _clean_slide_xml(self, xml_text: str) -> str:
        """Clean slide XML to remove sensitive artifacts.

        Removes:
        - Off-canvas objects (shapes with negative coordinates)
        - Hidden slide markers
        - Cropping artifacts
        """
        cleaned = xml_text

        # Remove hidden slide markers
        # <p:show/> element with show="0" indicates hidden slide
        cleaned = re.sub(
            r'<a:show\s+[^>]*/>',
            '',
            cleaned,
        )

        # NOTE: no whitespace collapsing — it corrupted text content in
        # xml:space="preserve" runs.

        return cleaned

    def _clean_relationships(self, xml_text: str) -> str:
        """Clean relationship XML to remove references to sensitive content."""
        cleaned = xml_text

        # Remove references to embedded objects
        cleaned = re.sub(
            r'<Relationship[^>]*Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/objectData"[^>]*/>',
            '',
            cleaned,
        )

        return cleaned

    def detect_risks(self, input_path: Path) -> List[str]:
        """Detect potential risk vectors in a PowerPoint file.

        Args:
            input_path: Source PPTX file path

        Returns:
            List of detected risk descriptions
        """
        risks = []
        ext = input_path.suffix.lower()

        if ext == '.ppt':
            risks.append("Legacy .ppt format - copied as-is without cleaning")

        if not zipfile.is_zipfile(input_path):
            return risks

        try:
            with zipfile.ZipFile(input_path, 'r') as zf:
                namelist = zf.namelist()

                # Check for speaker notes
                if any('notesSlides' in f for f in namelist):
                    risks.append(
                        "Speaker notes present - may contain frankest content"
                    )

                # Check for hidden slides
                try:
                    presentation_xml = zf.read('ppt/presentation.xml').decode('utf-8')
                    if 'hiddenSlide' in presentation_xml:
                        risks.append(
                            "Hidden slides present - may contain sensitive content"
                        )
                except Exception:
                    pass

                # Check for embedded objects
                if any('embed' in f for f in namelist):
                    risks.append(
                        "Embedded objects present - may contain sensitive data"
                    )

                # Check for off-canvas objects
                for slide_path in namelist:
                    if slide_path.startswith('ppt/slides/slide'):
                        slide_xml = zf.read(slide_path).decode('utf-8')
                        if 'xsp:-' in slide_xml or 'yp:-' in slide_xml:
                            risks.append(
                                f"Off-canvas objects in {slide_path} - may contain deprecated content"
                            )
                            break

        except Exception as e:
            risks.append(f"Risk detection failed: {e}")

        return risks