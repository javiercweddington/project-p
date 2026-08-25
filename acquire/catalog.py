"""
Catalog module - enriched file analysis built on top of the read module.

Provides a Catalog class that attaches to ProjectFile/ProjectGroup objects
from the read module and provides:
- File type catalog (images, videos, PDFs, documents, CAD files, etc.)
- File counts by type
- Sensitive information detection (company names, people, locations)

READ-ONLY: This module is strictly READ-ONLY. It only reads file metadata
and content for analysis purposes. No files are created, modified, or deleted.
"""

from __future__ import annotations

import re
import io
import os
import sys
import hashlib
import difflib
import logging
from pathlib import Path
from typing import Dict, List, Optional, Set, Tuple
from collections import defaultdict
from datetime import datetime
import zipfile

try:
    from pypdf import PdfReader
    HAS_PYPDF = True
except ImportError:
    HAS_PYPDF = False

try:
    from gliner import GLiNER
    HAS_GLINER = True
except ImportError:
    HAS_GLINER = False
    GLiNER = None

from .read import ProjectFile, ProjectGroup

# Import metadata classes for enhanced sensitivity detection
from .metadata import (
    ImageOCR,
    SensitiveDocFlagger,
    FilenamePatternDetector,
    CADMetadataExtractor,
)

# Module-level logger for proper error reporting
_logger = logging.getLogger(__name__)

# ---- GLiNER-based entity detection ----
# Based on working test: ~/test_glin.py
# Uses urchade/gliner_multi-v2.1 with labels ["person", "organization"]

_gliner_model = None
_gliner_model_name = os.environ.get("GLINER_MODEL", "urchade/gliner_multi-v2.1")
_gliner_threshold = float(os.environ.get("GLINER_THRESHOLD", "0.5"))
# Expanded labels based on cross-product analysis (Bosch + Wizama)
# GLiNER is zero-shot and supports arbitrary label names
_gliner_labels = [
    "person", "organization", "email", "phone", "address",
    "invoice", "date", "money",
]


def _get_gliner_model():
    """Get or create the GLiNER model (lazy singleton)."""
    global _gliner_model
    if _gliner_model is not None:
        return _gliner_model
    if not HAS_GLINER:
        return None
    try:
        _logger.info("Loading GLiNER model %s ...", _gliner_model_name)
        _gliner_model = GLiNER.from_pretrained(_gliner_model_name)
        _logger.info("GLiNER model loaded.")
        return _gliner_model
    except Exception as e:
        _logger.warning("Failed to load GLiNER model: %s", e)
        _gliner_model = object()  # sentinel to avoid retrying
        return None


def _extract_entities_with_gliner(text: str, source: str) -> List[EntityHit]:
    """Use GLiNER to extract named entities from text.

    Falls back to regex-based detection if GLiNER is not available.

    Args:
        text: Text content to scan.
        source: Source filename/path for attribution.

    Returns:
        List of EntityHit objects.
    """
    model = _get_gliner_model()
    if model is None:
        _logger.debug("GLiNER not available, falling back to regex-based entity detection")
        return extract_entities_from_text_regex(text, source)

    hits = []
    seen = set()  # deduplicate within the same text
    # Process text in chunks to handle long documents
    chunk_size = 3000
    overlap = 200
    for start in range(0, len(text), chunk_size - overlap):
        chunk = text[start:start + chunk_size]
        if not chunk.strip():
            continue
        try:
            for e in model.predict_entities(chunk, _gliner_labels, threshold=_gliner_threshold):
                entity_text = e['text'].strip()
                entity_type = e['label'].lower()
                confidence = round(float(e['score']), 3)
                # Deduplicate: keep highest confidence for each (type, value) pair
                key = (entity_type, entity_text.lower())
                if key in seen:
                    continue
                seen.add(key)
                # Map all GLiNER entity types, not just person/organization
                if entity_type == 'person':
                    hits.append(EntityHit(entity_type='person', value=entity_text,
                                          source=source, confidence=confidence))
                elif entity_type == 'organization':
                    hits.append(EntityHit(entity_type='company', value=entity_text,
                                          source=source, confidence=confidence))
                elif entity_type in ('email', 'phone', 'address', 'invoice', 'date', 'money'):
                    hits.append(EntityHit(entity_type=entity_type, value=entity_text,
                                          source=source, confidence=confidence))
        except Exception as e:
            _logger.debug("GLiNER prediction failed on chunk: %s", e)
            continue

    if hits:
        return hits

    # Fallback to regex if GLiNER found nothing
    return extract_entities_from_text_regex(text, source)


# ---- Sensitivity detection patterns ----

# Common company/organization name patterns (keywords that suggest a company)
_COMPANY_KEYWORDS = [
    'llc', 'inc', 'corp', 'corporation', 'ltd', 'gmbh', 'sa', 'plc',
    'company', 'co.', 'industries', 'technologies', 'systems', 'labs',
    'medical', 'solutions', 'group', 'holdings', 'partners',
]

# Title patterns that often precede person names
_PERSON_TITLE_PATTERN = re.compile(
    r'(?:Mr|Mrs|Ms|Dr|Prof|CEO|COO|CTO|CFO|VP|Director|Manager|Engineer)'
    r'[.\s]+([A-Z][a-z]+)\s+([A-Z][a-z]+)',
)

# Location patterns
_LOCATION_PATTERNS = [
    # US States
    r'\b(?:Alabama|Alaska|Arizona|Arkansas|California|Colorado|Connecticut|Delaware|Florida|Georgia|Hawaii|Idaho|Illinois|Indiana|Iowa|Kansas|Kentucky|Louisiana|Maine|Maryland|Massachusetts|Michigan|Minnesota|Mississippi|Missouri|Montana|Nebraska|Nevada|New Hampshire|New Jersey|New Mexico|New York|North Carolina|North Dakota|Ohio|Oklahoma|Oregon|Pennsylvania|Rhode Island|South Carolina|South Dakota|Tennessee|Texas|Utah|Vermont|Virginia|Washington|West Virginia|Wisconsin|Wyoming)\b',
    # Countries (English)
    r'\b(?:United States|USA|America|China|Japan|Germany|France|United Kingdom|UK|Canada|Australia|India|Brazil|Mexico|Italy|Spain|Korea|Hong Kong|Singapore|Switzerland|Netherlands|Sweden|Norway|Denmark|Finland|Belgium|Austria|Ireland|Portugal|Greece|Turkey|Russia|South Africa|Argentina|Chile|Colombia)\b',
    # Countries (Chinese)
    r'(?:中国|美国|日本|德国|法国|英国|加拿大|澳大利亚|印度|巴西|墨西哥|意大利|西班牙|韩国|香港|新加坡|瑞士|荷兰|瑞典|挪威|丹麦|芬兰|比利时|奥地利|爱尔兰|葡萄牙|希腊|土耳其|俄罗斯|南非|阿根廷|智利|哥伦比亚)',
    # Chinese Provinces
    r'(?:北京|上海|天津|重庆|河北|山西|辽宁|吉林|黑龙江|江苏|浙江|安徽|福建|江西|山东|河南|湖北|湖南|广东|海南|四川|贵州|云南|陕西|甘肃|青海|台湾|内蒙古|广西|西藏|宁夏|新疆|香港|澳门)',
    # Chinese Cities
    r'(?:北京|上海|广州|深圳|杭州|南京|成都|重庆|武汉|西安|长沙|郑州|济南|青岛|大连|沈阳|哈尔滨|长春|昆明|贵阳|南宁|福州|厦门|南昌|合肥|石家庄|太原|兰州|银川|西宁|乌鲁木齐|拉萨|呼和浩特|海口|三亚|东莞|佛山|苏州|无锡|常州|南通|徐州|扬州|镇江|泰州|盐城|淮安|宿迁|连云港|常州|温州|宁波|绍兴|嘉兴|湖州|金华|衢州|舟山|台州|丽水|杭州|南京|苏州|无锡|常州|徐州|南通|扬州|镇江|泰州|盐城|淮安|宿迁|连云港)',
    # Common city patterns (City, State or City, Country)
    r'\b[A-Z][a-z]+(?:\s+[A-Z][a-z]+)+,?\s+(?:[A-Z]{2}|\w+)\b',
]

# Address pattern (street addresses - English)
_ADDRESS_PATTERN_EN = re.compile(
    r'\b\d+\s+[A-Z][a-z]+(?:\s+[A-Z][a-z]+)*\s+(?:St|St|Ave|Blvd|Dr|Ln|Rd|Way|Pl|Ct|Ter|Cir)\.?\b',
    re.IGNORECASE,
)

# Address pattern (Chinese addresses)
_ADDRESS_PATTERN_CN = re.compile(
    r'(?:[^\s]{2,10}(?:省|市|区|县|镇|乡|村|路|街|大道|巷|弄|号|室|楼|栋|座|层|楼|单元))',
)

# ZIP/Postal code patterns (US)
_US_ZIP_PATTERN = re.compile(r'\b\d{5}(?:-\d{4})?\b')

# ZIP/Postal code patterns (China - 6 digits)
_CN_ZIP_PATTERN = re.compile(r'\b\d{6}\b')


class SensitivityFlag:
    """A flag indicating potentially sensitive information found in a file."""

    def __init__(self, flag_type: str, value: str, source: str,
                 context: str = "", confidence: float = 1.0):
        self.flag_type = flag_type  # 'company', 'person', 'location', 'email', 'phone'
        self.value = value          # The detected sensitive value
        self.source = source        # Filename or path where it was found
        self.context = context      # Surrounding text/context
        self.confidence = confidence  # 0.0 to 1.0

    def __repr__(self):
        return (f"SensitivityFlag(type={self.flag_type!r}, value={self.value!r}, "
                f"source={self.source!r})")


def extract_pdf_text(filepath: Path, max_pages: int = 5) -> str:
    """Extract text from a PDF file.

    Args:
        filepath: Path to the PDF file.
        max_pages: Maximum number of pages to extract (to limit processing).

    Returns:
        Extracted text as a string, or empty string if extraction fails.
    """
    if not HAS_PYPDF:
        return ""
    try:
        reader = PdfReader(str(filepath))
        pages = []
        for i in range(min(max_pages, len(reader.pages))):
            pages.append(reader.pages[i].extract_text() or "")
        return '\n'.join(pages)
    except Exception as e:
        _logger.debug("Failed to extract text from PDF %s: %s", filepath, e)
        return ""


def extract_pdf_text_from_zip(zip_path: Path, zip_entry: str, max_pages: int = 5) -> str:
    """Extract text from a PDF file inside a zip archive.

    Args:
        zip_path: Path to the zip file.
        zip_entry: Path of the PDF entry inside the zip.
        max_pages: Maximum number of pages to extract.

    Returns:
        Extracted text as a string, or empty string if extraction fails.
    """
    if not HAS_PYPDF:
        return ""
    try:
        with zipfile.ZipFile(zip_path, 'r') as zf:
            with zf.open(zip_entry) as pdf_file:
                reader = PdfReader(io.BytesIO(pdf_file.read()))
                pages = []
                for i in range(min(max_pages, len(reader.pages))):
                    pages.append(reader.pages[i].extract_text() or "")
                return '\n'.join(pages)
    except Exception as e:
        _logger.debug("Failed to extract text from PDF %s in zip %s: %s",
                       zip_entry, zip_path, e)
        return ""


def scan_text_for_sensitivity(text: str, source: str) -> List[SensitivityFlag]:
    """Scan extracted text for sensitive information.

    Args:
        text: Text content to scan.
        source: Source filename/path for attribution.

    Returns:
        List of SensitivityFlag objects for detected sensitive information.
    """
    flags = []
    text_lower = text.lower()

    # Email addresses
    for match in re.finditer(r'[\w\.-]+@[\w\.-]+\.\w+', text):
        flags.append(SensitivityFlag(
            flag_type='email',
            value=match.group(),
            source=source,
            context=_get_context(text, match.start()),
        ))

    # Phone numbers
    for match in re.finditer(
        r'(?:\+\d{1,3}[-.\s]\d{3,4}[-.\s]\d{3,8}|\(?\d{3}\)?[-.\s]\d{3,4}[-.\s]\d{4})',
        text
    ):
        flags.append(SensitivityFlag(
            flag_type='phone',
            value=match.group(),
            source=source,
            context=_get_context(text, match.start()),
        ))

    # Company names (look for patterns like "Company Name LLC" or "Name Technologies")
    # IMPORTANT: The regex requires the company keyword to be at the END of the match
    # to avoid false positives like "was machined in Kansas" matching "Kansas" as a company.
    for keyword in _COMPANY_KEYWORDS:
        pattern = re.compile(
            r'([A-Z][A-Za-z]+(?:\s+[A-Z][A-Za-z]+)*\s+' + re.escape(keyword) + r')',
            re.IGNORECASE,
        )
        for match in pattern.finditer(text):
            company = match.group(1).strip()
            if 2 <= len(company.split()) <= 5:  # Reasonable company name length
                flags.append(SensitivityFlag(
                    flag_type='company',
                    value=company,
                    source=source,
                    context=_get_context(text, match.start()),
                    confidence=0.6,
                ))

    # Person names via titles (e.g., "CEO Lech Murawski", "Dr. John Smith")
    for match in _PERSON_TITLE_PATTERN.finditer(text):
        first, last = match.group(1), match.group(2)
        name = f"{first} {last}"
        flags.append(SensitivityFlag(
            flag_type='person',
            value=name,
            source=source,
            context=_get_context(text, match.start()),
            confidence=0.8,
        ))

    # Locations (states, countries, Chinese provinces/cities)
    for pattern in _LOCATION_PATTERNS:
        for match in re.finditer(pattern, text):
            location = match.group().strip()
            if len(location) > 2:  # Avoid false positives
                flags.append(SensitivityFlag(
                    flag_type='location',
                    value=location,
                    source=source,
                    context=_get_context(text, match.start()),
                    confidence=0.7,
                ))

    # US ZIP codes (indicate addresses)
    for match in _US_ZIP_PATTERN.finditer(text):
        flags.append(SensitivityFlag(
            flag_type='address',
            value=match.group(),
            source=source,
            context=_get_context(text, match.start()),
            confidence=0.5,
        ))

    # Chinese ZIP codes (6 digits)
    for match in _CN_ZIP_PATTERN.finditer(text):
        flags.append(SensitivityFlag(
            flag_type='address',
            value=match.group(),
            source=source,
            context=_get_context(text, match.start()),
            confidence=0.4,  # Lower confidence - 6 digits could be other things
        ))

    # Chinese addresses
    for match in _ADDRESS_PATTERN_CN.finditer(text):
        flags.append(SensitivityFlag(
            flag_type='address',
            value=match.group(),
            source=source,
            context=_get_context(text, match.start()),
            confidence=0.6,
        ))

    return flags


def _get_context(text: str, pos: int, window: int = 50) -> str:
    """Get surrounding context text around a position."""
    start = max(0, pos - window)
    end = min(len(text), pos + window)
    context = text[start:end].replace('\n', ' ').strip()
    if start > 0:
        context = "..." + context
    if end < len(text):
        context = context + "..."
    return context


# File type classifications
IMAGE_EXTS = {
    '.jpg', '.jpeg', '.png', '.gif', '.bmp', '.svg', '.webp',
    '.tiff', '.tif', '.ico', '.heic', '.heif', '.raw', '.cr2',
}

VIDEO_EXTS = {
    '.mp4', '.avi', '.mov', '.mkv', '.wmv', '.flv', '.webm',
    '.m4v', '.mpg', '.mpeg', '.3gp',
}

AUDIO_EXTS = {
    '.mp3', '.wav', '.flac', '.aac', '.ogg', '.wma', '.m4a',
}

DOCUMENT_EXTS = {
    '.pdf', '.doc', '.docx', '.xls', '.xlsx', '.ppt', '.pptx',
    '.txt', '.rtf', '.odt', '.ods', '.odp', '.csv', '.tsv',
}

CAD_EXTS = {
    '.step', '.stp', '.stl', '.obj', '.fbx', '.blend',
    '.dwg', '.dxf', '.sldprt', '.sldasm', '.ipt', '.iam',
    '.prt', '.asm', '.x_t', '.x_b', '.iges', '.igs',
}

ELECTRONICS_EXTS = {
    '.brd', '.sch', '.pcb', '.kicad_pcb', '.kicad_sch',
    '.fzz', '.hex', '.bin', '.elf',
}

ARCHIVE_EXTS = {
    '.zip', '.tar', '.gz', '.bz2', '.7z', '.rar', '.xz',
}

FIRMWARE_EXTS = {
    '.hex', '.bin', '.elf', '.fw', '.uif', '.dfu',
}

CODE_EXTS = {
    '.py', '.js', '.c', '.cpp', '.h', '.java', '.go', '.rs',
    '.html', '.css', '.json', '.xml', '.yaml', '.yml', '.toml',
}


FILE_TYPE_MAP = {
    'image': IMAGE_EXTS,
    'video': VIDEO_EXTS,
    'audio': AUDIO_EXTS,
    'document': DOCUMENT_EXTS,
    'cad': CAD_EXTS,
    'electronics': ELECTRONICS_EXTS,
    'archive': ARCHIVE_EXTS,
    'firmware': FIRMWARE_EXTS,
    'code': CODE_EXTS,
}


def classify_file(filename: str) -> str:
    """Classify a file by extension into a type category."""
    ext = Path(filename).suffix.lower()
    for file_type, exts in FILE_TYPE_MAP.items():
        if ext in exts:
            return file_type
    return 'other'


class FileCatalog:
    """Catalog of file types within a single ProjectFile.

    Attaches to a ProjectFile and provides detailed file type analysis.
    """

    def __init__(self, project_file: ProjectFile):
        self.project_file = project_file
        self._files_by_type: Optional[Dict[str, List[str]]] = None
        self._file_count_by_type: Optional[Dict[str, int]] = None
        self._total_files: Optional[int] = None
        self._sensitivity_flags: Optional[List[SensitivityFlag]] = None

    @property
    def files_by_type(self) -> Dict[str, List[str]]:
        """Dictionary mapping file type to list of filenames."""
        if self._files_by_type is None:
            self._build_catalog()
        return self._files_by_type

    @property
    def file_count_by_type(self) -> Dict[str, int]:
        """Dictionary mapping file type to count of files."""
        if self._file_count_by_type is None:
            self._build_catalog()
        return self._file_count_by_type

    @property
    def total_files(self) -> int:
        if self._total_files is None:
            self._build_catalog()
        return self._total_files

    @property
    def has_images(self) -> bool:
        return self.file_count_by_type.get('image', 0) > 0

    @property
    def has_videos(self) -> bool:
        return self.file_count_by_type.get('video', 0) > 0

    @property
    def has_documents(self) -> bool:
        return self.file_count_by_type.get('document', 0) > 0

    @property
    def has_cad(self) -> bool:
        return self.file_count_by_type.get('cad', 0) > 0

    @property
    def has_audio(self) -> bool:
        return self.file_count_by_type.get('audio', 0) > 0

    @property
    def has_electronics(self) -> bool:
        return self.file_count_by_type.get('electronics', 0) > 0

    @property
    def has_firmware(self) -> bool:
        return self.file_count_by_type.get('firmware', 0) > 0

    def _build_catalog(self):
        """Build the file type catalog."""
        self._files_by_type = defaultdict(list)
        self._file_count_by_type = defaultdict(int)
        self._total_files = 0

        if self.project_file.is_zipped:
            self._catalog_zip()
        else:
            self._catalog_directory()

        # Convert defaultdicts to regular dicts
        self._files_by_type = dict(self._files_by_type)
        self._file_count_by_type = dict(self._file_count_by_type)

    def _catalog_zip(self):
        """Catalog files in a zip archive."""
        try:
            with zipfile.ZipFile(self.project_file.filepath, 'r') as zf:
                for name in zf.namelist():
                    if not name.endswith('/'):  # Skip directory entries
                        file_type = classify_file(name)
                        self._files_by_type[file_type].append(name)
                        self._file_count_by_type[file_type] += 1
                        self._total_files += 1
        except zipfile.BadZipFile as e:
            _logger.warning("Invalid zip file %s: %s", self.project_file.filepath, e)
        except PermissionError as e:
            _logger.warning("Permission denied reading zip %s: %s",
                            self.project_file.filepath, e)
        except OSError as e:
            _logger.warning("OS error reading zip %s: %s", self.project_file.filepath, e)

    def _catalog_directory(self):
        """Catalog files in a directory."""
        try:
            for filepath in self.project_file.filepath.rglob('*'):
                if filepath.is_file() and not filepath.name.startswith('._'):
                    # Get relative path
                    rel = str(filepath.relative_to(self.project_file.filepath))
                    file_type = classify_file(rel)
                    self._files_by_type[file_type].append(rel)
                    self._file_count_by_type[file_type] += 1
                    self._total_files += 1
        except PermissionError as e:
            _logger.warning("Permission denied reading directory %s: %s",
                            self.project_file.filepath, e)

    def get_sensitivity_flags(self, scan_pdf_content: bool = True) -> List[SensitivityFlag]:
        """Analyze filenames AND PDF content for sensitive information.

        Scans for patterns that indicate:
        - Company names
        - Person names
        - Email addresses
        - Phone numbers
        - Locations / addresses
        - Document types that may contain sensitive data

        Args:
            scan_pdf_content: If True, also extract and scan PDF text content.
        """
        if self._sensitivity_flags is not None:
            return self._sensitivity_flags

        self._sensitivity_flags = []

        # --- Phase 1: Scan filenames ---
        email_pattern = re.compile(r'[\w\.-]+@[\w\.-]+\.\w+')
        phone_pattern = re.compile(
            r'(?:'
            r'\+\d{1,3}[-.\s]\d{3,4}[-.\s]\d{3,8}'
            r'|\(?\d{3}\)?[-.\s]\d{3,4}[-.\s]\d{4}'
            r')'
        )
        sensitive_doc_types = [
            'nda', 'non-disclosure', 'confidential',
            'invoice', 'quotation', 'proposal', 'contract',
            'salary', 'payment', 'receipt',
        ]

        files_to_scan = []
        if self.project_file.is_zipped:
            try:
                with zipfile.ZipFile(self.project_file.filepath, 'r') as zf:
                    files_to_scan = zf.namelist()
            except zipfile.BadZipFile as e:
                _logger.warning("Invalid zip file %s: %s", self.project_file.filepath, e)
            except OSError as e:
                _logger.warning("OS error reading zip %s: %s", self.project_file.filepath, e)
        else:
            try:
                for filepath in self.project_file.filepath.rglob('*'):
                    if filepath.is_file():
                        files_to_scan.append(
                            str(filepath.relative_to(self.project_file.filepath))
                        )
            except PermissionError as e:
                _logger.warning("Permission denied reading directory %s: %s",
                                self.project_file.filepath, e)

        for filename in files_to_scan:
            if filename.endswith('/'):
                continue

            for match in email_pattern.finditer(filename):
                self._sensitivity_flags.append(SensitivityFlag(
                    flag_type='email', value=match.group(), source=filename,
                ))
            for match in phone_pattern.finditer(filename):
                self._sensitivity_flags.append(SensitivityFlag(
                    flag_type='phone', value=match.group(), source=filename,
                ))
            filename_lower = filename.lower()
            for doc_type in sensitive_doc_types:
                if doc_type in filename_lower:
                    self._sensitivity_flags.append(SensitivityFlag(
                        flag_type='sensitive_doc', value=doc_type,
                        source=filename, confidence=0.7,
                    ))
                    break

        # --- Phase 2: Scan PDF content ---
        if scan_pdf_content and HAS_PYPDF:
            pdf_files = []
            if self.project_file.is_zipped:
                try:
                    with zipfile.ZipFile(self.project_file.filepath, 'r') as zf:
                        pdf_files = [
                            n for n in zf.namelist()
                            if n.lower().endswith('.pdf') and not n.endswith('/')
                        ]
                except zipfile.BadZipFile as e:
                    _logger.warning("Invalid zip file %s: %s", self.project_file.filepath, e)
                except OSError as e:
                    _logger.warning("OS error reading zip %s: %s", self.project_file.filepath, e)
                for pdf_entry in pdf_files:
                    text = extract_pdf_text_from_zip(self.project_file.filepath, pdf_entry)
                    if text:
                        self._sensitivity_flags.extend(
                            scan_text_for_sensitivity(text, pdf_entry)
                        )
            else:
                try:
                    for filepath in self.project_file.filepath.rglob('*'):
                        if filepath.is_file() and filepath.suffix.lower() == '.pdf':
                            text = extract_pdf_text(filepath)
                            if text:
                                rel = str(filepath.relative_to(self.project_file.filepath))
                                self._sensitivity_flags.extend(
                                    scan_text_for_sensitivity(text, rel)
                                )
                except PermissionError as e:
                    _logger.warning("Permission denied reading directory %s: %s",
                                    self.project_file.filepath, e)

        return self._sensitivity_flags

    def summary(self) -> str:
        """Get a human-readable summary of the catalog."""
        lines = [
            f"File Catalog for: {self.project_file.project_name} (v{self.project_file.version})",
            f"Total files: {self.total_files}",
            "",
            "File types:",
        ]

        for file_type in sorted(self.file_count_by_type.keys()):
            count = self.file_count_by_type[file_type]
            lines.append(f"  {file_type}: {count}")

        if self._sensitivity_flags:
            lines.append("")
            lines.append(f"Sensitivity flags: {len(self._sensitivity_flags)}")
            # Count by type
            flag_counts = defaultdict(int)
            for flag in self._sensitivity_flags:
                flag_counts[flag.flag_type] += 1
            for flag_type, count in sorted(flag_counts.items()):
                lines.append(f"  {flag_type}: {count}")

        return '\n'.join(lines)


class ProjectCatalog:
    """Catalog for an entire ProjectGroup (multiple versions).

    Aggregates FileCatalog objects from all versions of a project.
    """

    def __init__(self, project_group: ProjectGroup):
        self.project_group = project_group
        self._version_catalogs: Dict[int, FileCatalog] = {}
        for f in project_group.files:
            self._version_catalogs[f.version] = FileCatalog(f)

    def get_version_catalog(self, version: int) -> Optional[FileCatalog]:
        """Get the catalog for a specific version."""
        return self._version_catalogs.get(version)

    @property
    def latest_catalog(self) -> Optional[FileCatalog]:
        """Get the catalog for the latest version."""
        if not self._version_catalogs:
            return None
        latest_version = max(self._version_catalogs.keys())
        return self._version_catalogs[latest_version]

    @property
    def all_sensitivity_flags(self) -> List[SensitivityFlag]:
        """Get all sensitivity flags across all versions."""
        flags = []
        for catalog in self._version_catalogs.values():
            flags.extend(catalog.get_sensitivity_flags())
        return flags

    def summary(self) -> str:
        """Get a human-readable summary of the project catalog."""
        lines = [
            f"Project Catalog: {self.project_group.project_name}",
            f"Company: {self.project_group.company}",
            f"Product: {self.project_group.product}",
            f"Versions: {self.project_group.version_count}",
            "",
        ]

        for version, catalog in sorted(self._version_catalogs.items()):
            lines.append(catalog.summary())
            lines.append("")

        return '\n'.join(lines)


def catalog_project(project_file: ProjectFile) -> FileCatalog:
    """Convenience function to create a FileCatalog for a ProjectFile."""
    return FileCatalog(project_file)


def catalog_project_group(project_group: ProjectGroup) -> ProjectCatalog:
    """Convenience function to create a ProjectCatalog for a ProjectGroup."""
    return ProjectCatalog(project_group)


# ---- Consolidation suggestions ----

class ConsolidationSuggestion:
    """A suggestion to consolidate two projects."""

    def __init__(self, project_a: ProjectGroup, project_b: ProjectGroup,
                 reason: str, confidence: float = 1.0):
        self.project_a = project_a
        self.project_b = project_b
        self.reason = reason
        self.confidence = confidence  # 0.0 to 1.0

    def __repr__(self):
        return (f"ConsolidationSuggestion({self.project_a.project_name!r} + "
                f"{self.project_b.project_name!r}, reason={self.reason!r})")


def _normalize_name(name: str) -> str:
    """Normalize a project name for comparison."""
    # Lowercase, remove special chars, collapse whitespace
    name = name.lower().strip()
    name = re.sub(r'[^a-z0-9\s-]', '', name)
    name = re.sub(r'\s+', ' ', name)
    return name


def suggest_consolidations(
    projects: List[ProjectGroup],
    name_threshold: float = 0.8
) -> List[ConsolidationSuggestion]:
    """Suggest projects that should be consolidated.

    Only suggests consolidation when:
    1. Same company (or very similar company name - typo case) AND same product
    2. Different project names (indicating version split or naming inconsistency)

    Args:
        projects: List of ProjectGroups to compare.
        name_threshold: Similarity threshold for company name typo detection (0.0 to 1.0).

    Returns:
        List of ConsolidationSuggestion objects.
    """
    suggestions = []
    compared = set()

    for i, a in enumerate(projects):
        for j, b in enumerate(projects):
            if i >= j:
                continue  # Only compare each pair once

            pair_key = (a.project_name, b.project_name)
            if pair_key in compared:
                continue
            compared.add(pair_key)

            # Skip if project names are identical (already same project)
            if a.project_name == b.project_name:
                continue

            # --- Check company match (exact or near-typo) ---
            company_exact = a.company.lower() == b.company.lower()
            company_similar = False
            if not company_exact:
                company_sim = difflib.SequenceMatcher(None, a.company.lower(), b.company.lower()).ratio()
                company_similar = company_sim >= name_threshold

            if not company_exact and not company_similar:
                continue  # Different companies - never suggest

            # --- Check product match (exact) ---
            if a.product.lower() != b.product.lower():
                continue  # Different products - never suggest

            # Same company (or typo) + same product + different project name
            if company_exact:
                reason = (f"Same company ({a.company}) + product ({a.product}) "
                          f"but different project names")
                confidence = 0.9
            else:
                reason = (f"Similar company name ({a.company} vs {b.company}, "
                          f"similarity: {company_sim:.2f}) + same product ({a.product})")
                confidence = company_sim * 0.9

            suggestions.append(ConsolidationSuggestion(
                project_a=a,
                project_b=b,
                reason=reason,
                confidence=confidence,
            ))

    return suggestions


def print_consolidation_suggestions(
    projects: List[ProjectGroup],
    name_threshold: float = 0.8
) -> None:
    """Print consolidation suggestions in a human-readable format."""
    suggestions = suggest_consolidations(projects, name_threshold)

    if not suggestions:
        print("No consolidation suggestions found.")
        return

    print(f"Consolidation Suggestions: {len(suggestions)}")
    print("=" * 60)

    for i, s in enumerate(suggestions, 1):
        print(f"\n{i}. {s.project_a.project_name} + {s.project_b.project_name}")
        print(f"   Reason: {s.reason}")
        print(f"   Confidence: {s.confidence:.0%}")
        print(f"   Company A: {s.project_a.company}, Product A: {s.project_a.product}")
        print(f"   Company B: {s.project_b.company}, Product B: {s.project_b.product}")
        print(f"   Versions A: {s.project_a.version_count}, "
              f"Versions B: {s.project_b.version_count}")


# ---- Named entity detection ----
#
# This module uses pattern-based detection rather than hardcoded name lists.
# It relies on structural cues (titles, capitalization, context words) to
# identify named entities in text.


class EntityHit:
    """A detected named entity in a file."""

    def __init__(self, entity_type: str, value: str, source: str,
                 confidence: float = 1.0, context: str = ""):
        self.entity_type = entity_type  # 'person', 'company', 'product'
        self.value = value
        self.source = source
        self.confidence = confidence
        self.context = context

    def __repr__(self):
        return (f"EntityHit(type={self.entity_type!r}, value={self.value!r}, "
                f"source={self.source!r})")


# Words that indicate a person name follows (titles, roles, positions)
_PERSON_TITLE_WORDS = {
    'mr', 'mrs', 'ms', 'miss', 'dr', 'prof', 'professor',
    'ceo', 'coo', 'cto', 'cfo', 'cmo', 'cpo',
    'vp', 'vice president', 'director', 'manager', 'engineer',
    'president', 'chairman', 'chairwoman', 'chief',
    'author', 'creator', 'designer', 'developer', 'writer',
    'contact', 'consultant', 'advisor', 'attorney', 'counsel',
    # Chinese titles
    '先生', '女士', '小姐', '博士', '教授', '经理', '总监', '总裁', '总经理',
}

# Words that indicate an organization/company context
_ORG_INDICATOR_WORDS = {
    'company', 'corporation', 'corporate', 'inc', 'incorporated', 'llc',
    'ltd', 'limited', 'gmbh', 'ag', 'sa', 'plc', 'nv', 'bv',
    'technologies', 'technology', 'industries', 'industry', 'solutions',
    'systems', 'labs', 'laboratories', 'medical', 'engineering',
    'group', 'groups', 'holdings', 'partners', 'partnership',
    'consulting', 'services', 'service', 'international', 'global',
}

# Product/model indicator patterns
_PRODUCT_PATTERNS = [
    # Model numbers: Model X-1234, Model 123, M-1234
    re.compile(r'(?:model|mod|mdl)[\s:-]*([A-Z0-9][-0-9A-Z]{2,10})', re.IGNORECASE),
    # Part numbers: P/N: 123, Part No: ABC-123
    re.compile(r'(?:part\s*(?:no|number)\.?|p/n|pn)[\s:]*([A-Z0-9][-0-9A-Z.]{2,20})', re.IGNORECASE),
    # Serial numbers: S/N: 123, Serial No: ABC
    re.compile(r'(?:serial\s*(?:no|number)\.?|s/n|sn)[\s:]*([A-Z0-9][-0-9A-Z.]{2,20})', re.IGNORECASE),
    # Versioned product names: ProductName v1.2, ProductName 1.2.3
    re.compile(r'([A-Z][A-Za-z0-9]+)\s+v(\d+\.\d+(?:\.\d+)?)', re.IGNORECASE),
]


# Simplified name validation - GLiNER handles entity detection,
# regex fallback uses minimal heuristics only.

def _is_likely_person_name(first: str, last: str) -> bool:
    """Heuristic check: does this look like a person name?

    Uses minimal structural heuristics:
    - Reasonable length (2-20 chars each)
    - Title Case format
    """
    if len(first) < 2 or len(last) < 2:
        return False
    if len(first) > 20 or len(last) > 20:
        return False
    # Check Title Case
    if not (first[0].isupper() and first[1:].islower()) and first.isalpha():
        return False
    if not (last[0].isupper() and last[1:].islower()) and last.isalpha():
        return False
    return True


def _detect_person_names_regex(text: str) -> List[str]:
    """Detect person names in text using regex patterns (fallback when LLM unavailable).


    Pattern 1: Title + Name (Dr. John Smith, CEO Jane Doe)
    Pattern 2: "Name, Title" (John Smith, CEO or John Smith, Manager)
    Pattern 3: Chinese title + single character surname
    Pattern 4: "by Name" pattern (common in documents)
    Pattern 5: "Name <email>" pattern
    """
    names = set()

    # Pattern 1: Title + First Last
    for match in re.finditer(
        r'(?:Mr|Mrs|Ms|Miss|Dr|Prof|Professor|CEO|COO|CTO|CFO|VP|Director|Manager|Engineer|President)[.:\s]+('
        r'[A-Z][a-z]+)\s+([A-Z][a-z]+)',
        text
    ):
        first, last = match.group(1), match.group(2)
        if _is_likely_person_name(first, last):
            names.add(f"{first} {last}")

    # Pattern 2: "Name, Title" (John Smith, CEO or John Smith, Manager)
    for match in re.finditer(
        r'([A-Z][a-z]+)\s+([A-Z][a-z]+),\s+(?:CEO|COO|CTO|CFO|VP|Director|Manager|Engineer|President|Owner|Founder)',
        text
    ):
        first, last = match.group(1), match.group(2)
        if _is_likely_person_name(first, last):
            names.add(f"{first} {last}")

    # Pattern 3: Chinese title + single character surname
    for match in re.finditer(
        r'(?:先生|女士|小姐|博士|教授|经理|总监|总裁|总经理)[\s:]*'
        r'([赵钱孙李周吴郑王冯陈褚卫蒋沈韩杨朱秦尤许何吕施张孔曹严华金魏陶姜'
        r'戚谢邹喻柏水窦章云苏潘葛奚范彭郎鲁韦昌马苗凤花方俞任袁柳酆鲍史唐费'
        r'廉岑薛雷贺倪汤滕殷罗毕郝邬安常乐于时傅皮卞齐康伍余元卜顾孟平黄和穆'
        r'萧尹姚邵湛汪祁毛禹狄米贝明臧计伏成戴谈宋茅庞熊纪舒屈项祝董梁杜阮蓝'
        r'闵席季麻强贾路娄危江童颜郭梅盛林刁钟徐邱骆高夏蔡田樊胡凌霍虞万支柯'
        r'昝管卢莫经房裘缪干解应宗丁宣贲邓郁单杭洪包诸左石崔吉钮龚程嵇邢滑裴'
        r'陆荣翁荀羊於惠甄曲家封芮羿储靳汲邴糜松井段富巫乌焦巴弓牧隗山谷车侯'
        r'宓蓬全郗班仰秋仲伊宫宁仇栾暴甘钭厉戎祖武符刘景詹束龙叶幸司韶郜黎蓟'
        r'薄印宿白怀蒲邰从鄂索咸籍赖卓蔺屠蒙池乔阴胥能苍双温庄晏瞿蛮充慕连茹'
        r'习宦艾鱼容向古易慎戈廖庾终暨居衡步都耿满弘匡国文寇广禄阙东欧殳沃利'
        r'蔚越夔隆师巩厍聂晁勾敖融冷訾辛阚那简饶空曾毋沙乜养鞠须丰巢关蒯相查'
        r'后荆红游竺权逯盖益桓公])',
        text
    ):
        names.add(match.group(1))

    # Pattern 4: "by Name" pattern (common in documents)
    for match in re.finditer(
        r'\bby\s+([A-Z][a-z]+)\s+([A-Z][a-z]+)\b',
        text
    ):
        first, last = match.group(1), match.group(2)
        if _is_likely_person_name(first, last):
            names.add(f"{first} {last}")

    # Pattern 5: "Name <email>" pattern
    for match in re.finditer(
        r'([A-Z][a-z]+)\s+([A-Z][a-z]+)\s+<[\w.+-]+@[\w.-]+>',
        text
    ):
        first, last = match.group(1), match.group(2)
        if _is_likely_person_name(first, last):
            names.add(f"{first} {last}")

    return list(names)


def _detect_company_names_regex(text: str) -> List[str]:
    """Detect company names in text using regex patterns.

    
    Pattern 1: Words followed by org indicators (X Technologies, Y LLC)
    Pattern 2: "Name & Name" pattern
    Pattern 3: "Name of Name" pattern
    
    IMPORTANT: The regex requires the org indicator word to be at the END
    of the match to avoid false positives like "was machined in Kansas".
    """
    companies = set()
    
    # Pattern 1: TitleCase words followed by org indicator (keyword at end)
    for keyword in _ORG_INDICATOR_WORDS:
        # Match: "Acme Technologies" but NOT "machined in Kansas"
        # The keyword must be the last word, preceded by TitleCase words
        pattern = re.compile(
            r'\b([A-Z][A-Za-z]+(?:\s+[A-Z][A-Za-z]+)*\s+' + re.escape(keyword) + r')\b',
            re.IGNORECASE,
        )
        for match in pattern.finditer(text):
            company = match.group(1).strip()
            words = company.split()
            if 2 <= len(words) <= 6:
                companies.add(company)
    
    # Pattern 2: "Name & Name" (partnerships, law firms)
    for match in re.finditer(
        r'([A-Z][A-Za-z]+)\s+&\s+([A-Z][A-Za-z]+)',
        text
    ):
        companies.add(match.group(0).strip())
    
    # Pattern 3: "Name of Name" (University of X, State of Y)
    for match in re.finditer(
        r'([A-Z][A-Za-z]+\s+of\s+[A-Z][A-Za-z]+)',
        text
    ):
        value = match.group(1).strip()
        if len(value.split()) >= 3:
            companies.add(value)
    
    return list(companies)


def _detect_product_names_regex(text: str) -> List[str]:
    """Detect product names and model numbers in text using regex patterns.

    
    Uses indicator patterns (Model, P/N, S/N, version patterns).
    """
    products = set()
    
    for pattern in _PRODUCT_PATTERNS:
        for match in pattern.finditer(text):
            value = match.group(0).strip()
            if len(value) > 2:
                products.add(value)
    
    return list(products)


def extract_entities_from_text_regex(text: str, source: str) -> List[EntityHit]:
    """Extract named entities from text content using regex patterns (fallback).

    Args:
        text: Text content to scan.
        source: Source filename/path for attribution.

    Returns:
        List of EntityHit objects for detected entities.
    """
    hits = []

    for name in _detect_person_names_regex(text):
        hits.append(EntityHit(
            entity_type='person',
            value=name,
            source=source,
            confidence=0.7,
        ))

    for company in _detect_company_names_regex(text):
        hits.append(EntityHit(
            entity_type='company',
            value=company,
            source=source,
            confidence=0.6,
        ))

    for product in _detect_product_names_regex(text):
        hits.append(EntityHit(
            entity_type='product',
            value=product,
            source=source,
            confidence=0.5,
        ))

    return hits


def extract_entities_from_text(text: str, source: str) -> List[EntityHit]:
    """Extract named entities from text content.

    Uses GLiNER when available, falls back to regex patterns.

    Args:
        text: Text content to scan.
        source: Source filename/path for attribution.

    Returns:
        List of EntityHit objects for detected entities.
    """
    return _extract_entities_with_gliner(text, source)


def extract_entities_from_filename(filename: str) -> List[EntityHit]:
    """Extract named entities from a filename.

    Args:
        filename: The filename to scan.

    Returns:
        List of EntityHit objects for detected entities.
    """
    hits = []

    # Person names in filename
    for name in _detect_person_names_regex(filename):
        hits.append(EntityHit(
            entity_type='person',
            value=name,
            source=filename,
            confidence=0.8,  # Higher confidence for filenames
        ))

    # Company names in filename
    for company in _detect_company_names_regex(filename):
        hits.append(EntityHit(
            entity_type='company',
            value=company,
            source=filename,
            confidence=0.7,
        ))

    # Product names in filename
    for product in _detect_product_names_regex(filename):
        hits.append(EntityHit(
            entity_type='product',
            value=product,
            source=filename,
            confidence=0.6,
        ))

    return hits


# ---- Batch file discrimination ----

class FileSignature:
    """A signature for a file used for comparison."""

    def __init__(self, filepath: Path, project_name: str, version: int,
                 filename: str, file_type: str, size: int,
                 content_hash: str = "", name_normalized: str = ""):
        self.filepath = filepath
        self.project_name = project_name
        self.version = version
        self.filename = filename
        self.file_type = file_type
        self.size = size
        self.content_hash = content_hash
        self.name_normalized = name_normalized or _normalize_name(filename)


class DuplicateGroup:
    """A group of files with identical content."""

    def __init__(self, content_hash: str):
        self.content_hash = content_hash
        self.files: List[FileSignature] = []

    def add(self, sig: FileSignature):
        self.files.append(sig)

    def __repr__(self):
        return (f"DuplicateGroup(hash={self.content_hash[:12]}..., "
                f"count={len(self.files)})")


class FileDiscriminator:
    """Compare files across projects to find duplicates, similarities, and unique files.

    Uses both content (hashes) and metadata (name, size, type) for discrimination.
    """

    def __init__(self, projects: List[ProjectGroup], sample_size: int = 8192,
                 scan_content: bool = False):
        """
        Args:
            projects: List of ProjectGroups to analyze.
            sample_size: Number of bytes to sample for hashing (for large files).
            scan_content: If True, also scan file contents for named entities.
        """
        self.projects = projects
        self.sample_size = sample_size
        self.scan_content = scan_content
        self._signatures: List[FileSignature] = []
        self._entities: Optional[List[EntityHit]] = None
        self._built = False

    def _build_signatures(self):
        """Build file signatures for all files in all projects."""
        if self._built:
            return

        self._signatures = []

        for pg in self.projects:
            for pf in pg.files:
                if pf.is_zipped:
                    self._signature_from_zip(pf)
                else:
                    self._signature_from_dir(pf)

        self._built = True

    def _signature_from_zip(self, pf: ProjectFile):
        """Build signatures for files in a zip archive."""
        try:
            with zipfile.ZipFile(pf.filepath, 'r') as zf:
                for name in zf.namelist():
                    if name.endswith('/'):
                        continue
                    file_type = classify_file(name)
                    info = zf.getinfo(name)
                    content_hash = self._hash_zip_entry(zf, name)
                    norm = _normalize_name(Path(name).stem)
                    self._signatures.append(FileSignature(
                        filepath=pf.filepath,
                        project_name=pf.project_name,
                        version=pf.version,
                        filename=name,
                        file_type=file_type,
                        size=info.file_size,
                        content_hash=content_hash,
                        name_normalized=norm,
                    ))
        except zipfile.BadZipFile as e:
            _logger.warning("Invalid zip file %s: %s", pf.filepath, e)
        except PermissionError as e:
            _logger.warning("Permission denied reading zip %s: %s", pf.filepath, e)
        except OSError as e:
            _logger.warning("OS error reading zip %s: %s", pf.filepath, e)

    def _signature_from_dir(self, pf: ProjectFile):
        """Build signatures for files in a directory."""
        try:
            for filepath in pf.filepath.rglob('*'):
                if not filepath.is_file() or filepath.name.startswith('._'):
                    continue
                file_type = classify_file(filepath.name)
                size = filepath.stat().st_size
                content_hash = self._hash_file(filepath)
                norm = _normalize_name(filepath.stem)
                self._signatures.append(FileSignature(
                    filepath=filepath,
                    project_name=pf.project_name,
                    version=pf.version,
                    filename=filepath.name,
                    file_type=file_type,
                    size=size,
                    content_hash=content_hash,
                    name_normalized=norm,
                ))
        except PermissionError as e:
            _logger.warning("Permission denied reading directory %s: %s",
                            pf.filepath, e)

    def _hash_file(self, filepath: Path) -> str:
        """Compute MD5 hash of a file (sample first N bytes for large files)."""
        h = hashlib.md5()
        try:
            with open(filepath, 'rb') as f:
                h.update(f.read(self.sample_size))
            return h.hexdigest()
        except PermissionError as e:
            _logger.debug("Permission denied reading file %s: %s", filepath, e)
            return ""
        except OSError as e:
            _logger.debug("OS error reading file %s: %s", filepath, e)
            return ""

    def _hash_zip_entry(self, zf: zipfile.ZipFile, name: str) -> str:
        """Compute MD5 hash of a zip entry (sample first N bytes)."""
        h = hashlib.md5()
        try:
            with zf.open(name) as f:
                h.update(f.read(self.sample_size))
            return h.hexdigest()
        except Exception as e:
            _logger.debug("Failed to hash zip entry %s: %s", name, e)
            return ""

    @property
    def signatures(self) -> List[FileSignature]:
        """Get all file signatures."""
        self._build_signatures()
        return self._signatures

    def find_duplicates(self) -> List[DuplicateGroup]:
        """Find files with identical content (same hash)."""
        self._build_signatures()
        hash_groups = defaultdict(list)
        for sig in self._signatures:
            if sig.content_hash:
                hash_groups[sig.content_hash].append(sig)

        return [
            DuplicateGroup(h) for h, sigs in hash_groups.items()
            if len(sigs) > 1
        ]

    def find_similar_names(self, threshold: float = 0.8) -> List[Tuple[FileSignature, FileSignature, float]]:
        """Find files with similar names but potentially different content."""
        self._build_signatures()
        similar = []
        compared = set()

        for i, a in enumerate(self._signatures):
            for j, b in enumerate(self._signatures):
                if i >= j:
                    continue
                pair = (id(a), id(b))
                if pair in compared:
                    continue
                compared.add(pair)

                sim = difflib.SequenceMatcher(None, a.name_normalized, b.name_normalized).ratio()
                if sim >= threshold and a.name_normalized != b.name_normalized:
                    similar.append((a, b, sim))

        return sorted(similar, key=lambda x: x[2], reverse=True)

    def find_by_type(self, file_type: str) -> List[FileSignature]:
        """Find all files of a given type."""
        self._build_signatures()
        return [s for s in self._signatures if s.file_type == file_type]

    def find_unique(self) -> List[FileSignature]:
        """Find files that are unique (no duplicate by hash, no similar name)."""
        self._build_signatures()
        dup_hashes = set()
        hash_groups = defaultdict(list)
        for sig in self._signatures:
            if sig.content_hash:
                hash_groups[sig.content_hash].append(sig)
        for h, sigs in hash_groups.items():
            if len(sigs) > 1:
                dup_hashes.add(h)

        similar_names = set()
        for i, a in enumerate(self._signatures):
            for j, b in enumerate(self._signatures):
                if i >= j:
                    continue
                sim = difflib.SequenceMatcher(None, a.name_normalized, b.name_normalized).ratio()
                if sim >= 0.8 and a.name_normalized != b.name_normalized:
                    similar_names.add(id(a))
                    similar_names.add(id(b))

        return [
            s for s in self._signatures
            if s.content_hash not in dup_hashes and id(s) not in similar_names
        ]

    def _extract_entities_from_signatures(self):
        """Extract named entities from all file signatures."""
        if self._entities is not None:
            return

        self._entities = []

        for sig in self._signatures:
            # Extract entities from filename
            self._entities.extend(extract_entities_from_filename(sig.filename))

        # Optionally scan file contents for entities (multi-format)
        if self.scan_content:
            for sig in self._signatures:
                text = self._extract_text_for_entity_scan(sig)
                if text:
                    self._entities.extend(extract_entities_from_text(text, sig.filename))

    def _extract_text_for_entity_scan(self, sig: FileSignature) -> str:
        """Extract text from a file for entity scanning (multi-format support)."""
        ext = sig.filename.lower().split('.')[-1] if '.' in sig.filename else ''
        filepath = sig.filepath

        # PDF
        if ext == 'pdf' and HAS_PYPDF:
            if filepath.suffix.lower() == '.zip':
                return extract_pdf_text_from_zip(filepath, sig.filename)
            else:
                return extract_pdf_text(filepath)

        # DOCX
        elif ext == 'docx':
            try:
                import docx
                doc = docx.Document(str(filepath))
                return '\n'.join(p.text for p in doc.paragraphs)[:50000]
            except Exception:
                return ""

        # XLSX
        elif ext == 'xlsx':
            try:
                import openpyxl
                wb = openpyxl.load_workbook(str(filepath), read_only=True, data_only=True)
                cells = []
                for ws in wb.worksheets:
                    for row in ws.iter_rows(values_only=True):
                        cells.append(' '.join(str(c) for c in row if c))
                wb.close()
                return '\n'.join(cells)[:50000]
            except Exception:
                return ""

        # PPTX
        elif ext == 'pptx':
            try:
                from pptx import Presentation
                prs = Presentation(str(filepath))
                texts = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            texts.append(shape.text)
                return '\n'.join(texts)[:50000]
            except Exception:
                return ""

        # Images - use OCR
        elif ext in ('png', 'jpg', 'jpeg', 'bmp', 'tiff', 'tif', 'webp'):
            try:
                ocr = ImageOCR()
                return ocr.extract_text(filepath) or ""
            except Exception:
                return ""

        # CAD (STEP, DXF) - text-based
        elif ext in ('step', 'stp', 'dxf'):
            try:
                return filepath.read_text(encoding='utf-8', errors='ignore')[:50000]
            except Exception:
                return ""

        # Plain text
        elif ext in ('txt', 'csv', 'md', 'log'):
            try:
                return filepath.read_text(encoding='utf-8', errors='ignore')[:50000]
            except Exception:
                return ""

        return ""

    @property
    def entities(self) -> List[EntityHit]:
        """Get all extracted named entities."""
        self._build_signatures()
        self._extract_entities_from_signatures()
        return self._entities

    def find_entities(self, entity_type: Optional[str] = None) -> List[EntityHit]:
        """Find named entities, optionally filtered by type.

        Args:
            entity_type: Filter by type ('person', 'company', 'product').
                        If None, return all entities.

        Returns:
            List of EntityHit objects.
        """
        if entity_type:
            return [e for e in self.entities if e.entity_type == entity_type]
        return self.entities

    def find_people(self) -> List[EntityHit]:
        """Find all person names detected."""
        return self.find_entities('person')

    def find_companies(self) -> List[EntityHit]:
        """Find all company names detected."""
        return self.find_entities('company')

    def find_products(self) -> List[EntityHit]:
        """Find all product names detected."""
        return self.find_entities('product')

    def summary(self) -> str:
        """Get a human-readable summary of file discrimination."""
        self._build_signatures()
        dups = self.find_duplicates()
        by_type = defaultdict(list)
        for s in self._signatures:
            by_type[s.file_type].append(s)

        lines = [
            "File Discrimination Summary",
            f"Projects analyzed: {len(self.projects)}",
            f"Total files: {len(self._signatures)}",
            "",
            "Files by type:",
        ]
        for file_type in sorted(by_type.keys()):
            sigs = by_type[file_type]
            lines.append(f"  {file_type}: {len(sigs)}")

        lines.append("")
        lines.append(f"Duplicate groups: {len(dups)}")
        for dup in dups:
            lines.append(f"  Hash {dup.content_hash[:12]}... ({len(dup.files)} files):")
            for f in dup.files:
                lines.append(f"    - {f.filename} ({f.project_name} v{f.version})")

        unique = self.find_unique()
        lines.append("")
        lines.append(f"Unique files: {len(unique)}")

        # Named entities
        all_entities = self.entities
        if all_entities:
            lines.append("")
            lines.append("Named Entities:")
            
            # Count by type
            entity_counts = defaultdict(int)
            for e in all_entities:
                entity_counts[e.entity_type] += 1
            for etype, count in sorted(entity_counts.items()):
                lines.append(f"  {etype}: {count}")
            
            lines.append("")
            lines.append("People detected:")
            people = self.find_people()
            if people:
                # Deduplicate and show
                seen = set()
                for e in people:
                    if e.value not in seen:
                        lines.append(f"  - {e.value} (from {e.source})")
                        seen.add(e.value)
            else:
                lines.append("  (none)")
            
            lines.append("")
            lines.append("Companies detected:")
            companies = self.find_companies()
            if companies:
                seen = set()
                for e in companies:
                    if e.value not in seen:
                        lines.append(f"  - {e.value} (from {e.source})")
                        seen.add(e.value)
            else:
                lines.append("  (none)")
            
            lines.append("")
            lines.append("Products detected:")
            products = self.find_products()
            if products:
                seen = set()
                for e in products:
                    if e.value not in seen:
                        lines.append(f"  - {e.value} (from {e.source})")
                        seen.add(e.value)
            else:
                lines.append("  (none)")

        return '\n'.join(lines)